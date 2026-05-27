from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = 10.0   # slide width in inches
H = 5.625  # slide height in inches

# MSO rectangle shape type ID
RECT = 1


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, x, y, w, h, color: str):
    shape = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _text(slide, txt: str, x, y, w, h, size=16, bold=False,
          color="FFFFFF", align=PP_ALIGN.LEFT, face="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = face
    return tb


def _bullets(slide, items: list, x, y, w, h, size=15, color="2D3748", face="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(color)
        run.font.name = face
    return tb


# ── Slide builders ────────────────────────────────────────────────────────────

def _title_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["darkBg"])

    # Left accent stripe
    _rect(slide, 0, 0, 0.35, H, t["accentColor"])

    # Bottom bar
    _rect(slide, 0, 4.5, W, 1.125, t["primaryColor"])

    # Title
    _text(slide, data.get("title", ""), 0.65, 1.3, 9.0, 1.4,
          size=42, bold=True, color="FFFFFF", align=PP_ALIGN.LEFT)

    # Subtitle
    sub = data.get("subtitle", "")
    if sub:
        _text(slide, sub, 0.65, 2.85, 8.5, 0.85,
              size=20, color=t["secondaryColor"], align=PP_ALIGN.LEFT)

    # Footer label
    _text(slide, "AI Generated • Claude", 0.65, 4.65, 5, 0.5,
          size=11, color="FFFFFF", align=PP_ALIGN.LEFT)


def _content_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["lightBg"])

    # Header bar
    _rect(slide, 0, 0, W, 1.05, t["primaryColor"])
    _text(slide, data.get("title", ""), 0.45, 0.12, 9.1, 0.82,
          size=26, bold=True, color="FFFFFF")

    # Left accent line
    _rect(slide, 0.45, 1.25, 0.07, 3.9, t["accentColor"])

    bullets = data.get("bullets", [])
    if bullets:
        _bullets(slide, bullets, 0.72, 1.28, 8.8, 4.0,
                 size=16, color="1E293B")


def _stats_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["lightBg"])

    # Header
    _rect(slide, 0, 0, W, 1.05, t["primaryColor"])
    _text(slide, data.get("title", ""), 0.45, 0.12, 9.1, 0.82,
          size=26, bold=True, color="FFFFFF")

    stats = data.get("stats", [])
    if not stats:
        return

    n = len(stats)
    gap = 0.25
    total_gap = gap * (n + 1)
    card_w = (W - total_gap) / n
    card_colors = [t["accentColor"], t["primaryColor"], t["accentColor"], t["primaryColor"]]

    for i, stat in enumerate(stats):
        cx = gap + i * (card_w + gap)
        _rect(slide, cx, 1.35, card_w, 3.3, card_colors[i % len(card_colors)])
        _text(slide, stat.get("value", ""), cx + 0.1, 1.65, card_w - 0.2, 1.5,
              size=38, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
        _text(slide, stat.get("label", ""), cx + 0.1, 3.2, card_w - 0.2, 1.1,
              size=13, color="FFFFFF", align=PP_ALIGN.CENTER)


def _two_column_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["lightBg"])

    # Header
    _rect(slide, 0, 0, W, 1.05, t["primaryColor"])
    _text(slide, data.get("title", ""), 0.45, 0.12, 9.1, 0.82,
          size=26, bold=True, color="FFFFFF")

    columns = (data.get("columns") or [])[:2]
    col_colors = [t["accentColor"], t["primaryColor"]]
    col_x = [0.35, 5.2]
    col_w = 4.55

    for i, col in enumerate(columns):
        cx = col_x[i]
        # Column header
        _rect(slide, cx, 1.2, col_w, 0.55, col_colors[i])
        _text(slide, col.get("title", ""), cx + 0.12, 1.25, col_w - 0.24, 0.45,
              size=15, bold=True, color="FFFFFF")
        # Column bullets
        _bullets(slide, col.get("bullets", []), cx + 0.12, 1.9, col_w - 0.15, 3.4,
                 size=13, color="1E293B")


def _section_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["primaryColor"])

    # Right accent block
    _rect(slide, 7.6, 0, 2.4, H, t["accentColor"])

    # Large title
    _text(slide, data.get("title", ""), 0.55, 1.6, 6.9, 2.2,
          size=36, bold=True, color="FFFFFF", align=PP_ALIGN.LEFT)

    # Accent underline bar
    _rect(slide, 0.55, 4.0, 2.2, 0.1, t["secondaryColor"])


def _conclusion_slide(prs, data: dict, t: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["darkBg"])

    # Left accent stripe
    _rect(slide, 0, 0, 0.35, H, t["accentColor"])

    # Bottom bar
    _rect(slide, 0, 4.5, W, 1.125, t["primaryColor"])

    # Title
    _text(slide, data.get("title", "Key Takeaways"), 0.65, 0.3, 9.0, 0.85,
          size=30, bold=True, color="FFFFFF")

    bullets = data.get("bullets", [])
    if bullets:
        _bullets(slide, bullets, 0.65, 1.3, 8.8, 3.0,
                 size=16, color=t["secondaryColor"])

    _text(slide, "Thank you", 0.65, 4.65, 4, 0.45,
          size=14, bold=True, color="FFFFFF")


# ── Public API ────────────────────────────────────────────────────────────────

_HANDLERS = {
    "title": _title_slide,
    "content": _content_slide,
    "stats": _stats_slide,
    "two_column": _two_column_slide,
    "section_break": _section_slide,
    "conclusion": _conclusion_slide,
}

_DEFAULT_THEME = {
    "primaryColor": "1E2761",
    "secondaryColor": "CADCFC",
    "accentColor": "F96167",
    "darkBg": "1E2761",
    "lightBg": "F5F7FA",
}


def generate_presentation(slide_json: dict, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    theme = {**_DEFAULT_THEME, **slide_json.get("theme", {})}

    for slide_data in slide_json.get("slides", []):
        slide_type = slide_data.get("type", "content")
        handler = _HANDLERS.get(slide_type, _content_slide)
        handler(prs, slide_data, theme)

    prs.save(output_path)
